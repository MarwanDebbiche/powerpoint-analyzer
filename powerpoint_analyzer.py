from pptx import Presentation
from collections import Counter


class PowerpointAnalyzer:

    def analyze_fonts(self, file):
        presentation = Presentation(file)

        font_data = {}

        for slide_idx, slide in enumerate(presentation.slides):
            shape_idx = 0
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                elif len(shape.text_frame.text):
                    for paragraph in shape.text_frame.paragraphs:
                        font_name = paragraph.font.name
                        print(f"PARAGRAPH FONT NAME : {font_name}")
                        if font_name is not None:
                            if font_name in font_data:
                                font_data[font_name]["count"] += 1
                                font_data[font_name]["text_data"].append(
                                    {
                                        "type": "paragraph",
                                        "text": paragraph.text,
                                        "slide": slide_idx + 1,
                                    }
                                )
                                font_data[font_name]["slides"].add(slide_idx + 1)

                            else:
                                font_data[font_name] = {}
                                font_data[font_name]["count"] = 1
                                font_data[font_name]["text_data"] = [
                                    {
                                        "type": "paragraph",
                                        "text": paragraph.text,
                                        "slide": slide_idx + 1,
                                    }
                                ]
                                font_data[font_name]["slides"] = {slide_idx + 1}
                        else:
                            for run in paragraph.runs:
                                font_name = run.font.name
                                print(f"RUN FONT NAME : {font_name}")
                                if font_name in font_data:
                                    font_data[font_name]["count"] += 1
                                    font_data[font_name]["text_data"].append(
                                        {
                                            "type": "run",
                                            "text": run.text,
                                            "slide": slide_idx + 1,
                                        }
                                    )
                                    font_data[font_name]["slides"].add(slide_idx + 1)

                                else:
                                    font_data[font_name] = {}
                                    font_data[font_name]["count"] = 1
                                    font_data[font_name]["text_data"] = [
                                        {
                                            "type": "run",
                                            "text": run.text,
                                            "slide": slide_idx + 1,
                                        }
                                    ]
                                    font_data[font_name]["slides"] = {slide_idx + 1}


                    shape_idx += 1

        return font_data
